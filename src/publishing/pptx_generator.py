from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import json
import requests
import tempfile
import os
from typing import List, Dict, Any, Tuple
from PIL import Image as PILImage

# Mapping of Archetype -> (Master Index, Layout Index)
LAYOUT_COORDS = {
    "hero": (0, 0),            # Master 0, Layout 0 ('Atmosphere')
    "documentary": (1, 2),     # Master 1, Layout 2 ('1 column')
    "split": (1, 4),           # Master 1, Layout 4 ('2 columns')
    "content_caption": (1, 3), # Master 1, Layout 3 ('1 column/photo')
    "grid": (1, 5),            # Master 1, Layout 5 ('3 columns')
    "table": (1, 2),           # Master 1, Layout 2 (Reuse 1 column)
    "blank": (2, 0)            # Master 2, Layout 0 ('Photo')
}

# Mapping of Archetype -> Placeholder Indices (based on 'Training Template.pptx')
PLACEHOLDER_MAP = {
    "hero": { "title": 0, "subtitle": 13 },
    "documentary": { "title": 0, "body": 13 },
    "split": { "title": 0, "left_body": 13, "right_body": 23 },
    "content_caption": { "title": 0, "body": 13, "image": 17 },
    "grid": { 
        "title": 0,
        "col1_img": 22, "col1_body": 13, 
        "col2_img": 23, "col2_body": 28,
        "col3_img": 24, "col3_body": 29
    },
    "blank": { "title": 0, "body": 13 }
}

def download_image(url: str, temp_dir: str) -> Tuple[str | None, Tuple[int, int] | None]:
    """
    Downloads an image from a URL and returns (local file path, (width, height)).
    """
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        
        # Determine file extension from content type or URL
        content_type = response.headers.get('content-type', '')
        if 'png' in content_type or '.png' in url.lower():
            ext = '.png'
        elif 'gif' in content_type or '.gif' in url.lower():
            ext = '.gif'
        else:
            ext = '.jpg'
        
        # Save to temp file
        filename = f"image_{hash(url) % 10000}{ext}"
        filepath = os.path.join(temp_dir, filename)
        
        with open(filepath, 'wb') as f:
            f.write(response.content)
        
        # Get image dimensions
        try:
            with PILImage.open(filepath) as img:
                dimensions = img.size
        except:
            dimensions = (400, 300)  # Default dimensions
        
        return filepath, dimensions
    except Exception as e:
        print(f"Failed to download image from {url}: {e}")
        return None, None

def parse_content_segments(markdown_text: str) -> List[Tuple[str, Any]]:
    """
    Parses markdown and returns a list of (type, content) tuples in order.
    Types: 'text', 'image'
    
    Alt text may contain size metadata in format: filename|{"bw":100,"bh":50,"cw":1280,"ch":720}
    where bw/bh = bbox width/height (image size on source), cw/ch = canvas width/height (source slide size)
    """
    if not markdown_text:
        return [('text', '')]
    
    # Pattern to match markdown images
    image_pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
    
    result = []
    last_end = 0
    
    for match in re.finditer(image_pattern, markdown_text):
        # Add text before this image (if any)
        text_before = markdown_text[last_end:match.start()].strip()
        if text_before:
            result.append(('text', text_before))
        
        # Parse alt text for size metadata
        alt_text = match.group(1)
        size_info = None
        filename = alt_text
        
        # Check if alt contains size metadata (format: filename|{json})
        if '|' in alt_text and '{' in alt_text:
            parts = alt_text.split('|', 1)
            filename = parts[0]
            try:
                size_info = json.loads(parts[1])
            except json.JSONDecodeError:
                pass
        
        # Add the image with parsed metadata
        result.append(('image', {
            'alt': filename,
            'url': match.group(2),
            'size_info': size_info  # May be None if not available
        }))
        
        last_end = match.end()
    
    # Add remaining text after last image (if any)
    remaining_text = markdown_text[last_end:].strip()
    if remaining_text:
        result.append(('text', remaining_text))
    
    return result if result else [('text', '')]

def estimate_text_height(text: str, width_inches: float = 8.5, font_size_pt: int = 11) -> float:
    """
    Estimate text height in inches based on content and font size.
    More conservative estimate to prevent overlap.
    """
    if not text:
        return 0
    
    lines = text.split('\n')
    line_count = 0
    
    for line in lines:
        line = line.strip()
        if not line:
            # Empty lines still take some space
            line_count += 0.5
            continue
        
        # Check for headers (take more vertical space)
        if line.startswith('#'):
            line_count += 1.5
            continue
            
        # Estimate characters that fit per line at given font size and width
        # At 11pt font, roughly 10 chars per inch
        chars_per_inch = 10
        chars_per_line = int(width_inches * chars_per_inch)
        
        # Calculate wrapped lines
        wrapped_lines = max(1, (len(line) // chars_per_line) + 1)
        line_count += wrapped_lines
    
    # Calculate height: ~0.25 inches per line at 11pt, plus padding
    line_height = 0.28 if font_size_pt <= 11 else 0.35
    base_height = line_count * line_height
    
    # Add padding for margins and safety
    return base_height + 0.3

def add_text_to_frame(text_frame, markdown_text):
    """
    Parses simple markdown (bullets, bold, headers) into a PPTX text frame.
    """
    if not markdown_text:
        return

    lines = markdown_text.split('\n')
    first_para = True
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Handle markdown headers - make them bold
        is_header = False
        if line.startswith('#'):
            line = re.sub(r'^#+\s*', '', line)
            is_header = True
            
        # Determine indent level (bullets)
        level = 0
        clean_line = line
        
        if line.startswith('- ') or line.startswith('* '):
            level = 0
            clean_line = line[2:]
        elif line.startswith('  - ') or line.startswith('  * '):
            level = 1
            clean_line = line[4:]
        
        if first_para:
            p = text_frame.paragraphs[0]
            first_para = False
        else:
            p = text_frame.add_paragraph()
        p.level = level
        
        # Handle Bold (**text**)
        parts = re.split(r'(\*\*.*?\*\*)', clean_line)
        
        for part in parts:
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            else:
                run.text = part
            run.font.size = Pt(11)
            if is_header:
                run.font.bold = True
                run.font.size = Pt(14)
                
        p.space_after = Pt(6)

def generate_pptx_document(project_title: str, nodes: List[Dict[str, Any]], output_path: str, template_path: str = None):
    """
    Generates a PPTX file from the project structure.
    Places text and images on the same slide using separate non-overlapping boxes.
    """
    
    def delete_all_slides(prs):
        """
        Robustly deletes all slides from the presentation by removing 
        both the slide ID entry and the relationship.
        """
        # Iterate backwards
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        try:
            delete_all_slides(prs)
        except Exception as e:
            print(f"Warning: Failed to clear template slides: {e}")
            # Continue anyway, better to have duplicates than crash or corruption
    else:
        prs = Presentation()
    
    # Slide dimensions (standard 10" x 7.5")
    SLIDE_WIDTH = 10.0
    SLIDE_HEIGHT = 7.5
    MARGIN = 0.5
    CONTENT_WIDTH = SLIDE_WIDTH - (2 * MARGIN)

    # Layout Mapping (Master Index, Layout Index)
    # Uses global LAYOUT_COORDS

    def get_layout(prs, master_idx, layout_idx):
        """Retrieves a specific layout from a specific master."""
        try:
            # Check if master exists
            if master_idx < len(prs.slide_masters):
                master = prs.slide_masters[master_idx]
                # Check if layout exists
                if layout_idx < len(master.slide_layouts):
                    return master.slide_layouts[layout_idx]
        except Exception as e:
            print(f"Error accessing layout m{master_idx}:l{layout_idx} - {e}")
        
        # Fallback to defaults if specific not found
        try:
            return prs.slide_layouts[layout_idx]
        except:
             return prs.slide_layouts[0] # Ultimate fallback

    def get_placeholder_safe(slide, idx):
        """Safely retrieve a placeholder by idx, or return None."""
        if idx is None: return None
        try:
            return slide.placeholders[idx]
        except KeyError:
            return None

    # 1. Title Slide
    # Use mapped hero layout
    hero_coords = LAYOUT_COORDS["hero"]
    title_slide_layout = get_layout(prs, hero_coords[0], hero_coords[1])
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Safety: check if placeholders exist
    if slide.shapes.title:
        slide.shapes.title.text = project_title
    
    # Try to find subtitle placeholder
    ph_indices = PLACEHOLDER_MAP["hero"]
    sub_ph = get_placeholder_safe(slide, ph_indices.get("subtitle"))
    if sub_ph:
        sub_ph.text = "Generated by Training Consolidation Workbench"
    
    # 2. Content Slides
    sorted_nodes = sorted(nodes, key=lambda x: x.get('order', 0))
    
    with tempfile.TemporaryDirectory() as temp_dir:
        for node in sorted_nodes:
            node_title = node.get('title', 'Untitled Section')
            content_md = node.get('content_markdown', '')
            target_layout = node.get('target_layout', 'documentary')
            
            # Map layout
            master_idx, layout_idx = LAYOUT_COORDS.get(target_layout, LAYOUT_COORDS['documentary'])
            layout = get_layout(prs, master_idx, layout_idx)
            
            # Get placeholder map for this layout
            ph_indices = PLACEHOLDER_MAP.get(target_layout, PLACEHOLDER_MAP['documentary'])

            # Parse into ordered segments
            segments = parse_content_segments(content_md)
            text_segments = [s for s in segments if s[0] == 'text']
            image_segments = [s for s in segments if s[0] == 'image']

            # Create slide
            slide = prs.slides.add_slide(layout)
            
            # Set Title
            title_idx = ph_indices.get("title", 0)
            if slide.shapes.title:
                slide.shapes.title.text = node_title
            elif title_idx is not None:
                title_ph = get_placeholder_safe(slide, title_idx)
                if title_ph: title_ph.text = node_title
            
            # --- Content Injection Strategy ---
            
            if target_layout == 'grid':
                # Grid Logic (3 Columns)
                cols = [
                    {'img': ph_indices.get('col1_img'), 'body': ph_indices.get('col1_body')},
                    {'img': ph_indices.get('col2_img'), 'body': ph_indices.get('col2_body')},
                    {'img': ph_indices.get('col3_img'), 'body': ph_indices.get('col3_body')}
                ]
                
                for i, col in enumerate(cols):
                    # Image
                    if i < len(image_segments) and col['img']:
                        ph = get_placeholder_safe(slide, col['img'])
                        if ph:
                            img_path, _ = download_image(image_segments[i][1]['url'], temp_dir)
                            if img_path:
                                ph.insert_picture(img_path)
                    
                    # Text
                    if i < len(text_segments) and col['body']:
                        ph = get_placeholder_safe(slide, col['body'])
                        if ph:
                            if not ph.has_text_frame: continue
                            add_text_to_frame(ph.text_frame, text_segments[i][1])

            elif target_layout == 'split':
                # Split Logic: Left=Text, Right=Image (or text 2)
                # Left
                left_ph = get_placeholder_safe(slide, ph_indices.get("left_body"))
                if left_ph and len(text_segments) > 0:
                    add_text_to_frame(left_ph.text_frame, text_segments[0][1])
                
                # Right
                right_ph = get_placeholder_safe(slide, ph_indices.get("right_body"))
                if right_ph:
                    if len(image_segments) > 0:
                        img_path, _ = download_image(image_segments[0][1]['url'], temp_dir)
                        if img_path:
                            try:
                                right_ph.insert_picture(img_path)
                            except AttributeError: pass
                    elif len(text_segments) > 1:
                        add_text_to_frame(right_ph.text_frame, text_segments[1][1])

            elif target_layout == 'content_caption':
                # Image + Caption
                img_ph = get_placeholder_safe(slide, ph_indices.get("image"))
                text_ph = get_placeholder_safe(slide, ph_indices.get("body"))
                
                if img_ph and len(image_segments) > 0:
                    img_path, _ = download_image(image_segments[0][1]['url'], temp_dir)
                    if img_path:
                        img_ph.insert_picture(img_path)
                
                if text_ph and len(text_segments) > 0:
                    add_text_to_frame(text_ph.text_frame, text_segments[0][1])

            else:
                # Documentary / Default
                body_ph = get_placeholder_safe(slide, ph_indices.get("body"))
                
                if body_ph:
                    # Combine all text
                    full_text = "\n\n".join([s[1] for s in text_segments])
                    add_text_to_frame(body_ph.text_frame, full_text)
                    
                    # TODO: Handle images in documentary mode if needed (create new slides?)
                else:
                    # Valid Manual Fallback Trigger
                    # If we really are 'blank' or failed to find main body
                    if target_layout == 'blank':
                        pass
            
                # Parse title (already done)
            
            # --- Manual Positioning Fallback (Only for BLANK or overflow) ---
            if target_layout == 'blank':
                # Use existing manual logic
                title_box = slide.shapes.add_textbox(
                    Inches(MARGIN), Inches(0.3), Inches(CONTENT_WIDTH), Inches(0.6)
                ) # ... (rest of manual title code)
            
            # For now, to keep the diff clean and functional, we will STOP here and let the loop continue
            # strictly for the manual parts IF layout was blank.
            # But since we replaced the 'BLANK' constant, we need to adapt the manual code below 
            # to only run if we didn't use placeholders effectively.
            
            # SIMPLIFICATION:
            # If we used a template layout (not blank), we skip manual positioning logic 
            # EXCEPT for 'blank' layout.
            
            if target_layout != 'blank':
                continue

            # --- Below is the manual positioning code for BLANK layout ---
            
            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(MARGIN), Inches(0.3), Inches(CONTENT_WIDTH), Inches(0.6)
            )
            title_tf = title_box.text_frame
            title_p = title_tf.paragraphs[0]
            title_p.text = node_title
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            
            # Track vertical position
            current_top = 1.0  # Start below title
            
            for seg_type, seg_content in segments:
                if seg_type == 'text' and seg_content:
                    # Estimate text height
                    text_height = estimate_text_height(seg_content, CONTENT_WIDTH)
                    text_height = max(0.5, min(text_height, SLIDE_HEIGHT - current_top - MARGIN))
                    
                    # Create text box
                    text_box = slide.shapes.add_textbox(
                        Inches(MARGIN), 
                        Inches(current_top), 
                        Inches(CONTENT_WIDTH), 
                        Inches(text_height)
                    )
                    text_box.text_frame.word_wrap = True
                    add_text_to_frame(text_box.text_frame, seg_content)
                    
                    # Add spacing after text (more generous to prevent overlap)
                    current_top += text_height + 0.3
                    
                elif seg_type == 'image':
                    # Download image
                    img_path, dimensions = download_image(seg_content['url'], temp_dir)
                    
                    if img_path:
                        try:
                            # Get size metadata if available (from alt text)
                            size_info = seg_content.get('size_info')
                            
                            # Destination slide dimensions in inches
                            DEST_WIDTH = SLIDE_WIDTH
                            DEST_HEIGHT = SLIDE_HEIGHT
                            
                            if size_info:
                                # Use actual source dimensions from metadata
                                # bw/bh = bbox width/height (image size on source)
                                # cw/ch = canvas width/height (source slide size)
                                bbox_width = size_info.get('bw', 0)
                                bbox_height = size_info.get('bh', 0)
                                canvas_width = size_info.get('cw', 1280)
                                canvas_height = size_info.get('ch', 720)
                                
                                if bbox_width > 0 and bbox_height > 0:
                                    # Calculate ratio of image to source canvas
                                    width_ratio = bbox_width / canvas_width
                                    height_ratio = bbox_height / canvas_height
                                    
                                    # Apply same ratio to destination slide
                                    img_width = DEST_WIDTH * width_ratio
                                    img_height = DEST_HEIGHT * height_ratio
                                else:
                                    # Fallback to pixel dimensions
                                    if dimensions:
                                        width_ratio = dimensions[0] / canvas_width
                                        height_ratio = dimensions[1] / canvas_height
                                        img_width = DEST_WIDTH * width_ratio
                                        img_height = DEST_HEIGHT * height_ratio
                                    else:
                                        img_width = 4
                                        img_height = 2.5
                            elif dimensions:
                                # No metadata - estimate from image pixel dimensions
                                SOURCE_CANVAS_WIDTH = 1280
                                SOURCE_CANVAS_HEIGHT = 720
                                img_px_width, img_px_height = dimensions
                                
                                width_ratio = img_px_width / SOURCE_CANVAS_WIDTH
                                height_ratio = img_px_height / SOURCE_CANVAS_HEIGHT
                                
                                img_width = DEST_WIDTH * width_ratio
                                img_height = DEST_HEIGHT * height_ratio
                            else:
                                # Complete fallback
                                img_width = 4
                                img_height = 2.5
                            
                            # Cap at reasonable maximums (don't exceed 90% of slide)
                            max_width = CONTENT_WIDTH * 0.9
                            max_height = 4.0  # Leave room for text
                            
                            if img_width > max_width:
                                scale = max_width / img_width
                                img_width *= scale
                                img_height *= scale
                                
                            if img_height > max_height:
                                scale = max_height / img_height
                                img_width *= scale
                                img_height *= scale
                            
                            # Ensure minimum size
                            img_width = max(img_width, 2.0)
                            img_height = max(img_height, 1.5)
                            # Center horizontally
                            left = (SLIDE_WIDTH - img_width) / 2
                            
                            slide.shapes.add_picture(
                                img_path, 
                                Inches(left), 
                                Inches(current_top), 
                                width=Inches(img_width)
                            )
                            
                            current_top += img_height + 0.3
                            
                        except Exception as e:
                            print(f"Failed to add image: {e}")
                
                # Check if we're running out of space - create continuation slide
                if current_top > SLIDE_HEIGHT - 1.0:
                    # Check if there are more segments
                    remaining = segments[segments.index((seg_type, seg_content)) + 1:]
                    if remaining:
                        slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
                        
                        # Add continuation title
                        title_box = slide.shapes.add_textbox(
                            Inches(MARGIN), Inches(0.3), Inches(CONTENT_WIDTH), Inches(0.6)
                        )
                        title_tf = title_box.text_frame
                        title_p = title_tf.paragraphs[0]
                        title_p.text = f"{node_title} (continued)"
                        title_p.font.size = Pt(28)
                        title_p.font.bold = True
                        
                        current_top = 1.0
                    
    prs.save(output_path)
    return output_path
