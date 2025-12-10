import streamlit as st
import yaml
from pptx import Presentation
import matplotlib.pyplot as plt
import matplotlib.patches as patches

st.set_page_config(layout="wide", page_title="PPTX Template Mapper")

def draw_layout_blueprint(layout):
    """
    Draws a visual representation of the slide layout using Matplotlib.
    This avoids needing complex PDF rendering libraries.
    """
    fig, ax = plt.subplots(figsize=(10, 6))
    
    # Slide dimensions (normalized to canvas)
    slide_width = 10
    slide_height = 5.625 # 16:9 aspect ratio approximation
    
    ax.set_xlim(0, slide_width)
    ax.set_ylim(0, slide_height)
    ax.invert_yaxis() # PPT coordinates start top-left
    
    # Draw Placeholders
    for shape in layout.placeholders:
        ph_idx = shape.placeholder_format.idx
        name = shape.name
        
        # Normalize coordinates
        x = shape.left.inches
        y = shape.top.inches
        w = shape.width.inches
        h = shape.height.inches
        
        # Create Rectangle
        rect = patches.Rectangle((x, y), w, h, linewidth=2, edgecolor='blue', facecolor='skyblue', alpha=0.3)
        ax.add_patch(rect)
        
        # Label with Index (Crucial for mapping)
        ax.text(x + w/2, y + h/2, f"ID: {ph_idx}\n({name})", 
                ha='center', va='center', fontsize=8, color='darkblue', weight='bold')

    ax.set_title(f"Layout: {layout.name}", fontsize=12)
    ax.axis('off')
    return fig

# --- APP LOGIC ---

st.title("🛠️ PowerPoint Template Mapper")
st.markdown("Upload your custom `.pptx` template to map the layouts to our 7 Archetypes.")

uploaded_file = st.file_uploader("Upload Template", type=["pptx"])

if uploaded_file:
    prs = Presentation(uploaded_file)
    
    # 1. Inspect all layouts
    layouts_flat = []
    layout_options = {} # "Master 0 - Layout 1: Name" -> (0, 1)
    
    for m_idx, master in enumerate(prs.slide_masters):
        for l_idx, layout in enumerate(master.slide_layouts):
            key = f"M{m_idx} / L{l_idx}: {layout.name}"
            layouts_flat.append(layout)
            layout_options[key] = (m_idx, l_idx, layout)

    # 2. Configuration Form
    archetypes = ["hero", "documentary", "split", "content_caption", "grid", "table", "blank"]
    
    # Store results here
    final_mapping = {"template_name": uploaded_file.name, "mappings": {}}

    col1, col2 = st.columns([1, 2])

    with col1:
        st.subheader("1. Archetype Mapping")
        selected_archetype = st.selectbox("Select Archetype to Map", archetypes)
        
        st.subheader("2. Select Matching Layout")
        selected_layout_name = st.selectbox(
            f"Which layout in the PPTX looks like a '{selected_archetype}'?",
            list(layout_options.keys())
        )
        
        m_idx, l_idx, current_layout = layout_options[selected_layout_name]
        
        # Show Blueprint
        st.pyplot(draw_layout_blueprint(current_layout))

    with col2:
        st.subheader("3. Map Placeholders")
        st.info(f"Look at the blueprint on the left. Match the 'ID' numbers to the roles below.")
        
        # Define required fields per archetype
        required_fields = {
            "hero": ["title", "subtitle"],
            "documentary": ["title", "body"],
            "split": ["title", "left_body", "right_body"],
            "content_caption": ["title", "image", "body"],
            "grid": ["title", "col1_img", "col1_body", "col2_img", "col2_body", "col3_img", "col3_body"],
            "table": ["title", "body"],
            "blank": ["title"]
        }
        
        fields = required_fields.get(selected_archetype, ["title"])
        
        ph_mapping = {}
        for field in fields:
            # Let user select from available IDs on this layout
            available_ids = [s.placeholder_format.idx for s in current_layout.placeholders]
            ph_mapping[field] = st.selectbox(f"Placeholder for '{field}'", available_ids, key=f"{selected_archetype}_{field}")

        # Save logic (in memory for display)
        final_mapping["mappings"][selected_archetype] = {
            "master_index": m_idx,
            "layout_index": l_idx,
            "placeholders": ph_mapping
        }

    # 4. Export
    st.divider()
    st.subheader("4. Generated YAML Config")
    st.code(yaml.dump(final_mapping, sort_keys=False), language="yaml")
    st.caption("Copy this into your `config/templates/` folder.")