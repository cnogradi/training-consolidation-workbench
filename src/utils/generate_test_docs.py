import os
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches

def generate_course_a_pdf(output_dir):
    filename = "course_a_intro_to_git.pdf"
    path = os.path.join(output_dir, filename)
    
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter
    
    # Page 1: Title
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 100, "Course A: Intro to Git")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, "Module 1: Git Basics")
    
    # Content
    c.drawString(100, height - 160, "Git is a distributed version control system.")
    c.drawString(100, height - 180, "Key concepts: Repository, Commit, Branch.")
    
    # Diagram: Commit Tree
    c.setStrokeColorRGB(0, 0, 0)
    c.setLineWidth(2)
    
    # Node 1
    c.circle(150, height - 300, 20, fill=0)
    c.drawString(140, height - 300, "C1")
    
    # Node 2
    c.circle(250, height - 300, 20, fill=0)
    c.drawString(240, height - 300, "C2")
    
    # Edge
    c.line(170, height - 300, 230, height - 300)
    
    c.drawString(100, height - 350, "Figure 1: A simple linear commit history.")
    
    c.showPage()
    c.save()
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Ray",
        "course_title": "Intro to Git",
        "version": "1.0",
        "scope_of_material": "Git Basics, Commits, Branches",
        "current_delivery_method": "Web",
        "duration_hours": 2.0,
        "costs": "Low",
        "tech_data_assessment": "None",
        "source_of_content": "Internal Wiki",
        "current_instructors": "Alice Smith",
        "audience": "Junior Developers",
        "location": "Remote",
        "level_of_material": "Beginner",
        "engineering_discipline": "Software",
        "comments": "Standard onboarding course."
    }
    return filename, path, metadata

def generate_course_b_pptx(output_dir):
    filename = "course_b_git_workflows.pptx"
    path = os.path.join(output_dir, filename)
    
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Course B: Git Workflows"
    subtitle.text = "Topic: Feature Branches and PRs"
    
    # Slide 2: Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Feature Branch Workflow"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "1. Create a branch from main"
    p = tf.add_paragraph()
    p.text = "2. Commit changes"
    p = tf.add_paragraph()
    p.text = "3. Open a Pull Request"
    
    # Slide 3: Diagram (Shapes)
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Merge Flowchart"
    
    # Main branch line
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(6), Inches(0.2))
    slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(0.5)).text_frame.text = "Main Branch"
    
    # Feature branch arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.CURVED_UP_ARROW, Inches(2), Inches(3.5), Inches(3), Inches(1))
    arrow.text_frame.text = "Feature Branch"
    
    prs.save(path)
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Collin",
        "course_title": "Git Workflows",
        "version": "2.1",
        "scope_of_material": "Feature Branching, Pull Requests",
        "current_delivery_method": "In Person",
        "duration_hours": 4.0,
        "costs": "Medium",
        "tech_data_assessment": "Required",
        "source_of_content": "DevOps Handbook",
        "current_instructors": "Bob Jones",
        "audience": "All Developers",
        "location": "New York Office",
        "level_of_material": "Intermediate",
        "engineering_discipline": "Software",
        "comments": "Critical for compliance."
    }
    return filename, path, metadata

def generate_course_c_docx(output_dir):
    filename = "course_c_git_collaboration.docx"
    path = os.path.join(output_dir, filename)
    
    doc = Document()
    doc.add_heading('Course C: Collaborating with Git', 0)
    
    doc.add_heading('Resolving Conflicts', level=1)
    doc.add_paragraph('Conflicts occur when two branches modify the same line of code.')
    
    doc.add_heading('Merge vs Rebase', level=1)
    table = doc.add_table(rows=3, cols=2)
    table.style = 'Table Grid'
    
    # Header
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Strategy'
    hdr_cells[1].text = 'Pros/Cons'
    
    # Row 1
    row1 = table.rows[1].cells
    row1[0].text = 'Merge'
    row1[1].text = 'Preserves history, but can be messy.'
    
    # Row 2
    row2 = table.rows[2].cells
    row2[0].text = 'Rebase'
    row2[1].text = 'Clean history, but rewrites commits.'
    
    doc.save(path)
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Pat",
        "course_title": "Collaborating with Git",
        "version": "1.5",
        "scope_of_material": "Conflicts, Merge vs Rebase",
        "current_delivery_method": "Hybrid",
        "duration_hours": 3.0,
        "costs": "Low",
        "tech_data_assessment": "None",
        "source_of_content": "Community Guidelines",
        "current_instructors": "Charlie Brown",
        "audience": "Senior Developers",
        "location": "London Office",
        "level_of_material": "Advanced",
        "engineering_discipline": "Software",
        "comments": "Focus on conflict resolution."
    }
    return filename, path, metadata

def generate_course_d_pdf(output_dir):
    """Mechanical Engineering - CAD Basics"""
    filename = "course_d_cad_basics.pdf"
    path = os.path.join(output_dir, filename)
    
    c = canvas.Canvas(path, pagesize=letter)
    width, height = letter
    
    # Page 1: Title
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 100, "Course D: CAD Basics")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, "Module 1: Introduction to 3D Modeling")
    
    # Content
    c.drawString(100, height - 160, "Computer-Aided Design (CAD) is essential for modern engineering.")
    c.drawString(100, height - 180, "Key concepts: Sketches, Constraints, Extrusions, Assemblies.")
    c.drawString(100, height - 200, "Common tools: SolidWorks, AutoCAD, Fusion 360.")
    
    # Simple diagram: 3D cube wireframe
    c.setStrokeColorRGB(0, 0, 0)
    c.setLineWidth(2)
    
    # Front face
    c.rect(150, height - 350, 100, 100, fill=0)
    # Back face (offset)
    c.rect(180, height - 320, 100, 100, fill=0)
    # Connecting lines
    c.line(150, height - 350, 180, height - 320)
    c.line(250, height - 350, 280, height - 320)
    c.line(250, height - 250, 280, height - 220)
    c.line(150, height - 250, 180, height - 220)
    
    c.drawString(100, height - 400, "Figure 1: Basic 3D wireframe representation.")
    
    c.showPage()
    c.save()
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Engineering Academy",
        "course_title": "CAD Basics",
        "version": "1.0",
        "scope_of_material": "3D Modeling, Sketches, Constraints, Extrusions",
        "current_delivery_method": "In Person",
        "duration_hours": 8.0,
        "costs": "High",
        "tech_data_assessment": "Required",
        "source_of_content": "ASME Guidelines",
        "current_instructors": "Dr. Jane Wilson",
        "audience": "Mechanical Engineers",
        "location": "Detroit Office",
        "level_of_material": "Beginner",
        "engineering_discipline": "Mechanical",
        "comments": "Hands-on lab required."
    }
    return filename, path, metadata

def generate_course_e_pptx(output_dir):
    """Mechanical Engineering - Thermodynamics"""
    filename = "course_e_thermodynamics.pptx"
    path = os.path.join(output_dir, filename)
    
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Course E: Thermodynamics I"
    subtitle.text = "Topic: Laws of Thermodynamics"
    
    # Slide 2: First Law
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "First Law of Thermodynamics"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Energy cannot be created or destroyed"
    p = tf.add_paragraph()
    p.text = "ΔU = Q - W"
    p = tf.add_paragraph()
    p.text = "Where: ΔU = Change in internal energy"
    p = tf.add_paragraph()
    p.text = "Q = Heat added to system, W = Work done by system"
    
    # Slide 3: Second Law
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Second Law of Thermodynamics"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Entropy of an isolated system always increases"
    p = tf.add_paragraph()
    p.text = "ΔS ≥ 0 for irreversible processes"
    p = tf.add_paragraph()
    p.text = "Applications: Heat engines, refrigeration cycles"
    
    # Slide 4: Carnot Cycle
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Carnot Cycle"
    
    # Add a simple diagram shape
    slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(2), Inches(3), Inches(2))
    textbox = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(2), Inches(0.5))
    textbox.text_frame.text = "Ideal Heat Engine"
    
    prs.save(path)
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Engineering Academy",
        "course_title": "Thermodynamics I",
        "version": "3.0",
        "scope_of_material": "Laws of Thermodynamics, Heat Engines, Entropy",
        "current_delivery_method": "Web",
        "duration_hours": 12.0,
        "costs": "Medium",
        "tech_data_assessment": "None",
        "source_of_content": "University Textbook",
        "current_instructors": "Prof. Robert Chen",
        "audience": "Engineering Students",
        "location": "Remote",
        "level_of_material": "Intermediate",
        "engineering_discipline": "Mechanical",
        "comments": "Core curriculum course."
    }
    return filename, path, metadata

def generate_course_f_docx(output_dir):
    """Mechanical Engineering - Materials Science"""
    filename = "course_f_materials_science.docx"
    path = os.path.join(output_dir, filename)
    
    doc = Document()
    doc.add_heading('Course F: Materials Science', 0)
    
    doc.add_heading('Material Properties', level=1)
    doc.add_paragraph('Understanding material properties is crucial for engineering design.')
    doc.add_paragraph('Key properties: Strength, Hardness, Ductility, Toughness, Elasticity.')
    
    doc.add_heading('Material Selection', level=1)
    doc.add_paragraph('Selection depends on application requirements and cost constraints.')
    
    # Add a table comparing materials
    table = doc.add_table(rows=4, cols=3)
    table.style = 'Table Grid'
    
    # Header
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Material'
    hdr_cells[1].text = 'Strength (MPa)'
    hdr_cells[2].text = 'Applications'
    
    # Row 1 - Steel
    row1 = table.rows[1].cells
    row1[0].text = 'Steel'
    row1[1].text = '400-2000'
    row1[2].text = 'Structures, automotive'
    
    # Row 2 - Aluminum
    row2 = table.rows[2].cells
    row2[0].text = 'Aluminum'
    row2[1].text = '100-600'
    row2[2].text = 'Aerospace, packaging'
    
    # Row 3 - Polymer
    row3 = table.rows[3].cells
    row3[0].text = 'Polymer'
    row3[1].text = '20-100'
    row3[2].text = 'Consumer products'
    
    doc.add_heading('Failure Modes', level=1)
    doc.add_paragraph('Common failure modes include:')
    doc.add_paragraph('• Fracture - brittle or ductile', style='List Bullet')
    doc.add_paragraph('• Fatigue - cyclic loading', style='List Bullet')
    doc.add_paragraph('• Creep - time-dependent deformation', style='List Bullet')
    doc.add_paragraph('• Corrosion - chemical degradation', style='List Bullet')
    
    doc.save(path)
    print(f"Generated {path}")
    
    metadata = {
        "business_unit": "Engineering Academy",
        "course_title": "Materials Science",
        "version": "2.0",
        "scope_of_material": "Material Properties, Selection, Failure Modes",
        "current_delivery_method": "Hybrid",
        "duration_hours": 10.0,
        "costs": "Medium",
        "tech_data_assessment": "Required",
        "source_of_content": "Industry Standards",
        "current_instructors": "Dr. Sarah Martinez",
        "audience": "Mechanical Engineers",
        "location": "Seattle Office",
        "level_of_material": "Advanced",
        "engineering_discipline": "Mechanical",
        "comments": "Lab sessions include tensile testing."
    }
    return filename, path, metadata

def generate_all(output_dir="test_docs"):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    docs = []
    # Software Engineering courses
    docs.append(generate_course_a_pdf(output_dir))
    docs.append(generate_course_b_pptx(output_dir))
    docs.append(generate_course_c_docx(output_dir))
    
    # Mechanical Engineering courses
    docs.append(generate_course_d_pdf(output_dir))
    docs.append(generate_course_e_pptx(output_dir))
    docs.append(generate_course_f_docx(output_dir))
    
    return docs

if __name__ == "__main__":
    generate_all()
